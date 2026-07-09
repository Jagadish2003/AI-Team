"""
R18-A1 / T2 — PowerPoint (.pptx) text extraction handler.

Extracts slide text from a presentation with ``python-pptx``: the text of every
shape that carries a text frame (titles, bullets, text boxes) plus table cell
text, grouped per slide with a ``# Slide N`` header. An encrypted deck (OLE
compound file, not a ZIP) is a loud
:data:`~discovery.ingest.extraction.ENCRYPTED` skip (AC3); a corrupt/unparseable
one raises :class:`~discovery.ingest.extraction.ExtractionError` so the ingestor
isolates it per-file (AC5). Speaker notes and embedded media are out of phase-one
scope.
"""
from __future__ import annotations

import io
import logging

from . import (
    ENCRYPTED,
    NO_HANDLER,
    ExtractedText,
    ExtractionError,
    ExtractionSkipped,
    ExtractionOutcome,
    register_handler,
)
from ._office_common import looks_encrypted_ooxml

logger = logging.getLogger(__name__)


def extract_pptx(raw: bytes, fmt: str) -> ExtractionOutcome:
    """Extract slide text from a .pptx, or record a loud skip."""
    if looks_encrypted_ooxml(raw):
        return ExtractionSkipped(ENCRYPTED, "PowerPoint deck is password-protected")

    try:
        from pptx import Presentation  # python-pptx
    except ImportError:  # pragma: no cover - python-pptx ships in requirements
        return ExtractionSkipped(NO_HANDLER, "python-pptx is not installed")

    try:
        presentation = Presentation(io.BytesIO(raw))
    except Exception as exc:  # noqa: BLE001 — a parse failure is a corrupt file (AC5)
        raise ExtractionError(f"could not parse .pptx: {type(exc).__name__}") from exc

    slides = list(presentation.slides)
    parts = []
    for index, slide in enumerate(slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            slide_parts.extend(_shape_text(shape))
        if slide_parts:
            parts.append(f"# Slide {index}")
            parts.extend(slide_parts)

    content = "\n".join(parts).strip()
    return ExtractedText(
        content=content,
        chunk_content_type="prose",
        structure_hints={"format": "pptx", "slides": len(slides), "chars": len(content)},
    )


def _shape_text(shape) -> list:
    """Return the non-empty text lines a slide shape carries (text frame or table)."""
    lines = []
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    lines.append("\t".join(cells))
    except Exception:  # noqa: BLE001 — one odd shape must not lose the whole slide
        pass
    return lines


register_handler("pptx", extract_pptx)
