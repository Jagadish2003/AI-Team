"""
R18-A1 / T2 (AT-524) — contract tests for the document format handlers.

Covers the acceptance criteria assigned to this subtask:

  AC1 — A text-based PDF, a .docx, an .xlsx, and a .pptx each yield extracted text
        with correct provenance (the ExtractedText a handler returns; delivery to
        the substrate + retrievability is the T3/R18-B1 concern).
  AC3 — A scanned-image PDF or an encrypted file is recorded as skipped-with-reason
        (ExtractionSkipped), never silently ingested as empty text.

The plain-text / markdown / CSV handler and format-detection routing are covered in
``test_document_extraction.py``. Corrupt-file isolation (raising ExtractionError)
is asserted here at the handler boundary and end-to-end in
``test_documents_ingestor.py`` (AC5).

Documents are built in-memory with the same parser libraries the handlers use, so
the tests are self-contained (no committed binary fixtures).
"""
from __future__ import annotations

import io

import pytest

from discovery.ingest import extraction
from discovery.ingest.documents import DocumentIngestor
from discovery.ingest.documents_source import DocumentRef, DocumentSource
from discovery.ingest.extraction import ExtractedText, ExtractionError, ExtractionSkipped


# ─────────────────────────────────────────────────────────────────────────────
# Document builders (real files, in memory)
# ─────────────────────────────────────────────────────────────────────────────
def _make_pdf(text: str, *, with_text: bool = True) -> bytes:
    """Build a minimal, valid single-page PDF (with a correct xref table).

    ``with_text=False`` produces a page with no text-showing operators — a stand-in
    for a scanned-image PDF (pages exist, no extractable text layer).
    """
    stream = (
        b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
        if with_text
        else b" "
    )
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objs, start=1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + obj + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (
        b"trailer\n<< /Size " + str(len(objs) + 1).encode() + b" /Root 1 0 R >>\n"
        b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    )
    return bytes(out)


def _make_encrypted_pdf(password: str = "secret") -> bytes:
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(_make_pdf("secret content")))
    writer = pypdf.PdfWriter()
    writer.append(reader)
    writer.encrypt(password)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _make_docx(paragraphs, *, author=None, title=None) -> bytes:
    import docx

    document = docx.Document()
    for para in paragraphs:
        document.add_paragraph(para)
    if author:
        document.core_properties.author = author
    if title:
        document.core_properties.title = title
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_xlsx(sheet_name, rows) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(title, body) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    box.text_frame.text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


# ─────────────────────────────────────────────────────────────────────────────
# AC1 — each format yields extracted text with correct provenance
# ─────────────────────────────────────────────────────────────────────────────
def test_ac1_pdf_text_extracted():
    out = extraction.extract(_make_pdf("Hello PDF world"), filename="report.pdf")
    assert isinstance(out, ExtractedText)
    assert "Hello PDF world" in out.content
    assert out.chunk_content_type == "prose"
    assert out.structure_hints["format"] == "pdf"
    assert out.structure_hints["pages"] == 1


def test_ac1_docx_text_and_provenance_extracted():
    raw = _make_docx(
        ["Quarterly revenue grew.", "Costs held flat."],
        author="Finance",
        title="Q2 Review",
    )
    out = extraction.extract(raw, filename="q2.docx")
    assert isinstance(out, ExtractedText)
    assert "Quarterly revenue grew." in out.content
    assert "Costs held flat." in out.content
    assert out.provenance["author"] == "Finance"
    assert out.provenance["title"] == "Q2 Review"


def test_ac1_docx_table_cells_included():
    import docx

    document = docx.Document()
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "metric"
    table.rows[0].cells[1].text = "value"
    table.rows[1].cells[0].text = "ARR"
    table.rows[1].cells[1].text = "1.2M"
    buf = io.BytesIO()
    document.save(buf)
    out = extraction.extract(buf.getvalue(), filename="t.docx")
    assert isinstance(out, ExtractedText)
    assert "metric\tvalue" in out.content
    assert "ARR\t1.2M" in out.content


def test_ac1_xlsx_sheet_and_cell_text_extracted():
    raw = _make_xlsx("Q2", [["metric", "value"], ["ARR", 123], ["NRR", "118%"]])
    out = extraction.extract(raw, filename="metrics.xlsx")
    assert isinstance(out, ExtractedText)
    assert "# Sheet: Q2" in out.content
    assert "metric\tvalue" in out.content
    assert "ARR\t123" in out.content
    assert out.structure_hints["sheets"] == ["Q2"]


def test_ac1_pptx_slide_text_extracted():
    raw = _make_pptx("Roadmap", "Ship document ingestion")
    out = extraction.extract(raw, filename="deck.pptx")
    assert isinstance(out, ExtractedText)
    assert "Roadmap" in out.content
    assert "Ship document ingestion" in out.content
    assert out.structure_hints["slides"] == 1


# ─────────────────────────────────────────────────────────────────────────────
# AC3 — scanned / encrypted files are loud skips, never silent emptiness
# ─────────────────────────────────────────────────────────────────────────────
def test_ac3_scanned_image_pdf_is_skipped_not_empty():
    out = extraction.extract(_make_pdf("", with_text=False), filename="scan.pdf")
    assert isinstance(out, ExtractionSkipped)
    assert out.reason == extraction.SCANNED_IMAGE


def test_ac3_encrypted_pdf_is_skipped():
    out = extraction.extract(_make_encrypted_pdf(), filename="locked.pdf")
    assert isinstance(out, ExtractionSkipped)
    assert out.reason == extraction.ENCRYPTED


@pytest.mark.parametrize("name", ["locked.docx", "locked.xlsx", "locked.pptx"])
def test_ac3_encrypted_ooxml_is_skipped(name):
    # An encrypted Office file is delivered as an OLE compound file, not a ZIP.
    out = extraction.extract(_OLE_MAGIC + b"encrypted-payload", filename=name)
    assert isinstance(out, ExtractionSkipped)
    assert out.reason == extraction.ENCRYPTED


# ─────────────────────────────────────────────────────────────────────────────
# Corrupt files raise (per-file isolation happens in the ingestor — AC5)
# ─────────────────────────────────────────────────────────────────────────────
def test_corrupt_pdf_raises_extraction_error():
    with pytest.raises(ExtractionError):
        extraction.extract(b"%PDF-1.4 not really a pdf at all", filename="bad.pdf")


def test_corrupt_docx_raises_extraction_error():
    # ZIP magic (so it is not treated as encrypted) but not a valid OOXML package.
    with pytest.raises(ExtractionError):
        extraction.extract(b"PK\x03\x04totally-broken", filename="bad.docx")


# ─────────────────────────────────────────────────────────────────────────────
# Registration — every phase-one format now has a handler
# ─────────────────────────────────────────────────────────────────────────────
@pytest.mark.parametrize("fmt", ["pdf", "docx", "xlsx", "pptx", "text", "markdown", "csv"])
def test_every_phase_one_format_has_a_handler(fmt):
    assert fmt in extraction._HANDLERS


# ─────────────────────────────────────────────────────────────────────────────
# End-to-end through the ingestor (T2 handler + T1 ingestor together)
# ─────────────────────────────────────────────────────────────────────────────
class _BytesSource(DocumentSource):
    def __init__(self, ref: DocumentRef, raw: bytes):
        self._ref = ref
        self._raw = raw

    def list_documents(self, org_id):
        return [self._ref]

    def read(self, org_id, ref):
        return self._raw


def test_docx_flows_through_ingestor_as_extracted_content():
    ref = DocumentRef(
        artifact_id="reports/q2.docx",
        filename="q2.docx",
        location="reports",
        signature="v1",
        source_timestamp="2026-06-11T08:00:00Z",
    )
    raw = _make_docx(["Revenue is up.", "Retention is strong."])
    ing = DocumentIngestor(source=_BytesSource(ref, raw))
    records = [r for b in ing.ingest_changes("org1", None) for r in b.records]
    assert len(records) == 1
    rec = records[0]
    assert rec["extraction"]["status"] == "extracted"
    assert rec["document_format"] == "docx"
    assert "Revenue is up." in rec["content"]
    assert rec["evidence_pointer"]["origin"] == "observed"
