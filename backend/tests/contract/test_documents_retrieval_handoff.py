"""
R18-A1 / T3 (AT-525) — document → retrieval hand-off, end to end (AC1, AC6).

Proves the acceptance criteria assigned to T3 against the REAL R18-B1 substrate,
through the one public pipeline a real caller uses:

    ingest_documents(org, source=…)        # T3 hand-off (this story)
      → retrieval.ingest_content(org, …)    # B1 T5 producer contract
        → embedder.embed_pending_for_org()  # B1 T3 async, gateway-only
          → retrieve(org, query)            # B1 T4 source-aware API

  AC1 — a text-based PDF, a .docx, an .xlsx, and a .pptx each yield extracted
        text with correct provenance, delivered to the substrate and subsequently
        retrievable.
  AC6 — the delivered content carries origin='observed' and the full
        EvidencePointer spine, and retrieval of that content shows the correct
        source file.

Embedding runs through a FAKE provider registered with the REAL model gateway and
selected via ``MODEL_EMBEDDING_PROVIDER`` (the production path executes; no direct
provider call), exactly as the B1 acceptance suite does. Provider/name are unique
to this module to avoid registry collisions.

The checkpoint side is driven through an in-memory store so the test exercises the
retrieval hand-off against the real DB without touching the checkpoint table.
"""
from __future__ import annotations

import io
import json
from typing import List

import pytest

from app import db
from app.model_gateway import register_provider
from app.model_gateway._interface import (
    GenerationRequest,
    GenerationResult,
    ModelProvider,
)
from app.provenance import OBSERVED
from app.retrieval import embedder
from app.retrieval.api import retrieve
from discovery.ingest.base import Checkpoint
from discovery.ingest.documents_handoff import ingest_documents
from discovery.ingest.documents_source import DocumentRef, DocumentSource


# ---------------------------------------------------------------------------
# Skip cleanly where there is no pgvector-backed store (mirrors the B1 suite).
# ---------------------------------------------------------------------------
def _retrieval_store_available() -> bool:
    try:
        con = db.connect()
        try:
            cur = con.cursor()
            cur.execute("SELECT to_regclass('public.retrieval_chunks')")
            return cur.fetchone()[0] is not None
        finally:
            con.close()
    except Exception:
        return False


pytestmark = pytest.mark.skipif(
    not _retrieval_store_available(),
    reason="retrieval_chunks store (pgvector) not present in this environment",
)

# The phase-one binary format handlers depend on these parsers (in requirements.txt;
# present in CI). Skip the module cleanly rather than error if a lean environment
# lacks one — the fast, parser-free wiring proof lives in
# discovery/tests/test_documents_handoff.py.
pytest.importorskip("pypdf")
pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("pptx")


# ---------------------------------------------------------------------------
# Fake embedding provider — content-shaped so ranking is real (one slot per term)
# ---------------------------------------------------------------------------
_TERMS = ("alpha", "beta", "gamma", "delta")


class _DocFakeProvider(ModelProvider):
    emits_own_telemetry = True

    def __init__(self, name: str, identity):
        self.name = name
        self._identity = identity

    def generate(self, req: GenerationRequest) -> GenerationResult:  # pragma: no cover
        return GenerationResult(text=None, provider=self.name, ok=False)

    def embed(self, texts: List[str]) -> List[List[float]]:
        out = []
        for t in texts:
            low = t.lower()
            out.append([1.0 if term in low else 0.0 for term in _TERMS] + [0.01])
        return out

    def embedding_identity(self):
        return self._identity


_DOC_PROVIDER = _DocFakeProvider("doc_handoff_embed", ("doc-handoff:model", "1"))
register_provider(_DOC_PROVIDER)


# ---------------------------------------------------------------------------
# Real document files (in memory), one per phase-one binary/text format.
# ---------------------------------------------------------------------------
def _make_pdf(text: str) -> bytes:
    stream = b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objs, start=1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + obj + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (
        b"trailer\n<< /Size " + str(len(objs) + 1).encode() + b" /Root 1 0 R >>\n"
        b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    )
    return bytes(out)


def _make_docx(text: str) -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph(text)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_xlsx(text: str) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["note", text])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Deck"
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# Each file's text carries one distinct term so retrieval ranking is unambiguous.
_FILES = {
    "lib/report.pdf": ("report.pdf", "application/pdf", _make_pdf("alpha quarterly review")),
    "lib/notes.docx": ("notes.docx", None, _make_docx("beta meeting notes and actions")),
    "lib/budget.xlsx": ("budget.xlsx", None, _make_xlsx("gamma budget figures")),
    "lib/deck.pptx": ("deck.pptx", None, _make_pptx("delta launch plan")),
}
_TERM_FOR = {
    "lib/report.pdf": "alpha",
    "lib/notes.docx": "beta",
    "lib/budget.xlsx": "gamma",
    "lib/deck.pptx": "delta",
}


class _FourFormatSource(DocumentSource):
    """A document source serving one real file of each phase-one format."""

    reports_deletes = True

    def list_documents(self, org_id):
        refs = []
        for artifact_id, (filename, content_type, _bytes) in _FILES.items():
            refs.append(
                DocumentRef(
                    artifact_id=artifact_id,
                    filename=filename,
                    location="lib",
                    signature=f"sig-{artifact_id}-1",
                    source_timestamp="2026-07-09T09:00:00Z",
                    content_type=content_type,
                    provenance={"library": "Corporate"},
                )
            )
        return refs

    def read(self, org_id, ref):
        return _FILES[ref.artifact_id][2]


class _Store:
    def __init__(self):
        self.data: dict = {}

    def read(self, org_id, connector_id):
        return self.data.get((org_id, connector_id))

    def save(self, cp: Checkpoint):
        self.data[(cp.org_id, cp.connector_id)] = cp


def _cleanup(org_id: str) -> None:
    con = db.connect()
    try:
        cur = con.cursor()
        cur.execute("DELETE FROM retrieval_chunks WHERE org_id = %s", (org_id,))
        con.commit()
    finally:
        con.close()


def _rows_for(org_id: str, source_artifact: str) -> list:
    con = db.connect()
    try:
        cur = con.cursor()
        cur.execute(
            "SELECT source_system, source_artifact, source_timestamp, provenance "
            "FROM retrieval_chunks WHERE org_id = %s AND source_artifact = %s",
            (org_id, source_artifact),
        )
        return [dict(r) for r in cur.fetchall()]
    finally:
        con.close()


@pytest.fixture
def org(request, monkeypatch):
    monkeypatch.setenv("MODEL_EMBEDDING_PROVIDER", _DOC_PROVIDER.name)
    name = f"doc_ho_{request.node.name}"[:60]
    _cleanup(name)
    yield name
    _cleanup(name)


# ===========================================================================
# AC1 — each format delivered to the substrate and subsequently retrievable
# ===========================================================================
def test_ac1_pdf_docx_xlsx_pptx_delivered_and_retrievable(org):
    store = _Store()
    result = ingest_documents(
        org,
        source=_FourFormatSource(),
        read_checkpoint=store.read,
        save_checkpoint=store.save,
    )

    # All four extracted formats were handed to the substrate and indexed.
    assert result.ok, result.error
    assert result.artifacts_handed_off == 4
    assert result.artifacts_indexed == 4
    assert result.artifacts_failed == 0

    # Each file produced at least one indexed chunk under the 'document' system.
    for artifact_id in _FILES:
        rows = _rows_for(org, artifact_id)
        assert rows, f"no chunks indexed for {artifact_id}"
        assert all(r["source_system"] == "document" for r in rows)

    # Embedding is async + gateway-driven; once it runs, each format's distinct
    # term retrieves its OWN file (subsequently retrievable — AC1).
    run = embedder.embed_pending_for_org(org)
    assert run.embedded == result.chunks_indexed
    for artifact_id, term in _TERM_FOR.items():
        hits = retrieve(org, term, k=5, source_filter=["document"])
        assert hits, f"{term!r} retrieved nothing for {artifact_id}"
        assert hits[0].source_artifact == artifact_id


# ===========================================================================
# AC6 — observed provenance + correct source file, on the stored + retrieved sides
# ===========================================================================
def test_ac6_observed_provenance_and_correct_source_file(org):
    store = _Store()
    ingest_documents(
        org,
        source=_FourFormatSource(),
        read_checkpoint=store.read,
        save_checkpoint=store.save,
    )

    # Persisted side: every chunk carries origin='observed', the EvidencePointer
    # spine, and the human-facing source file identity (AC6).
    rows = _rows_for(org, "lib/report.pdf")
    assert rows
    prov = json.loads(rows[0]["provenance"])
    assert prov["origin"] == OBSERVED
    assert prov["filename"] == "report.pdf"
    assert prov["evidence_pointer"]["origin"] == OBSERVED
    assert prov["evidence_pointer"]["source_artifact"] == "lib/report.pdf"

    # Retrieved side: the hit names the correct source file, and its evidence
    # pointer is observed (retrieved content was seen directly).
    embedder.embed_pending_for_org(org)
    hits = retrieve(org, "alpha", k=5, source_filter=["document"])
    assert hits
    top = hits[0]
    assert top.source_system == "document"
    assert top.source_artifact == "lib/report.pdf"
    assert top.to_evidence_pointer().origin == OBSERVED
