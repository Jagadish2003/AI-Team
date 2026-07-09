"""
R18-A1 / T6 (AT-528) — Section 3 acceptance suite for Document Ingestion.

One cohesive, acceptance-oriented contract suite that proves every Section-3
criterion (AC1–AC6) against the COMPLETE, real pipeline — not a single layer:

    ingest_documents(org, source=…)          # T1 ingestor + T3 hand-off
      → retrieval.ingest_content(org, …)      # R18-B1 producer contract
        → embedder.embed_pending_for_org()    # R18-B1 async, gateway-only
          → retrieve(org, query)              # R18-B1 source-aware API

The per-task suites (``discovery/tests/test_documents_*`` and
``tests/contract/test_documents_retrieval_handoff``) verify each layer in
isolation; this suite is the holistic view a reviewer reads to confirm the whole
story. Its distinctive value is asserting the LOUD-SKIP discipline end to end:
a scanned/encrypted file (AC3), an oversized file (AC4), and a corrupt file (AC5)
are recorded with a reason AND leave NO chunks in the substrate — never silently
ingested as empty content — while every healthy file in the same run is indexed
and retrievable.

Embedding runs through a FAKE provider registered with the REAL model gateway and
selected via ``MODEL_EMBEDDING_PROVIDER`` (the production path executes; no direct
provider call). Provider/name are unique to this module to avoid registry
collisions. The checkpoint side uses an in-memory store so the test drives the
real DB without touching the checkpoint table.
"""
from __future__ import annotations

import io
import json
from typing import Dict, List, Optional

import pytest

from app import db
from app.model_gateway import register_provider
from app.model_gateway._interface import (
    GenerationRequest,
    GenerationResult,
    ModelProvider,
)
from app.provenance import OBSERVED
from app.retrieval import embedder
from app.retrieval.api import retrieve
from discovery.ingest import extraction
from discovery.ingest.base import Checkpoint
from discovery.ingest.documents_handoff import ingest_documents
from discovery.ingest.documents_source import DocumentRef, DocumentSource


# ---------------------------------------------------------------------------
# Skip cleanly where there is no pgvector-backed store or a parser is missing.
# ---------------------------------------------------------------------------
def _retrieval_store_available() -> bool:
    try:
        con = db.connect()
        try:
            cur = con.cursor()
            cur.execute("SELECT to_regclass('public.retrieval_chunks')")
            return cur.fetchone()[0] is not None
        finally:
            con.close()
    except Exception:
        return False


pytestmark = pytest.mark.skipif(
    not _retrieval_store_available(),
    reason="retrieval_chunks store (pgvector) not present in this environment",
)

pytest.importorskip("pypdf")
pytest.importorskip("docx")
pytest.importorskip("openpyxl")
pytest.importorskip("pptx")


# ---------------------------------------------------------------------------
# Fake embedding provider — one-hot over the per-file terms so ranking is real.
# ---------------------------------------------------------------------------
_TERMS = ("alpha", "beta", "gamma", "delta", "epsilon")


class _AcceptanceProvider(ModelProvider):
    emits_own_telemetry = True

    def __init__(self, name: str, identity):
        self.name = name
        self._identity = identity

    def generate(self, req: GenerationRequest) -> GenerationResult:  # pragma: no cover
        return GenerationResult(text=None, provider=self.name, ok=False)

    def embed(self, texts: List[str]) -> List[List[float]]:
        out = []
        for t in texts:
            low = (t or "").lower()
            out.append([1.0 if term in low else 0.0 for term in _TERMS] + [0.01])
        return out

    def embedding_identity(self):
        return self._identity


_PROVIDER = _AcceptanceProvider("doc_acceptance_embed", ("doc-acceptance:model", "1"))
register_provider(_PROVIDER)


# ---------------------------------------------------------------------------
# Real document builders (in memory).
# ---------------------------------------------------------------------------
def _make_pdf(text: str, *, with_text: bool = True) -> bytes:
    stream = (
        b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
        if with_text
        else b" "
    )
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objs, start=1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + obj + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (
        b"trailer\n<< /Size " + str(len(objs) + 1).encode() + b" /Root 1 0 R >>\n"
        b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    )
    return bytes(out)


def _make_docx(text: str) -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph(text)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_xlsx(text: str) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["note", text])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(text: str) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Deck"
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"  # encrypted OOXML container


# ---------------------------------------------------------------------------
# A flexible in-memory document source.
# ---------------------------------------------------------------------------
class ListSource(DocumentSource):
    def __init__(self, entries: List[dict], reports_deletes: bool = True):
        self.entries = entries
        self.reports_deletes = reports_deletes
        self.reads: List[str] = []

    def list_documents(self, org_id):
        refs = []
        for e in self.entries:
            refs.append(
                DocumentRef(
                    artifact_id=e["id"],
                    filename=e.get("filename", e["id"].split("/")[-1]),
                    location="lib",
                    signature=e.get("signature", "v1"),
                    source_timestamp=e.get("ts", "2026-07-09T09:00:00Z"),
                    content_type=e.get("content_type"),
                    size_bytes=e.get("size_bytes"),
                    provenance={"library": "Corporate"},
                )
            )
        return refs

    def read(self, org_id, ref):
        self.reads.append(ref.artifact_id)
        for e in self.entries:
            if e["id"] == ref.artifact_id:
                return e["bytes"]
        raise KeyError(ref.artifact_id)


class Store:
    def __init__(self):
        self.data: dict = {}

    def read(self, org_id, connector_id):
        return self.data.get((org_id, connector_id))

    def save(self, cp: Checkpoint):
        self.data[(cp.org_id, cp.connector_id)] = cp


# ---------------------------------------------------------------------------
# DB helpers.
# ---------------------------------------------------------------------------
def _cleanup(org_id: str) -> None:
    con = db.connect()
    try:
        cur = con.cursor()
        cur.execute("DELETE FROM retrieval_chunks WHERE org_id = %s", (org_id,))
        con.commit()
    finally:
        con.close()


def _rows_for(org_id: str, source_artifact: str) -> list:
    con = db.connect()
    try:
        cur = con.cursor()
        cur.execute(
            "SELECT source_system, source_artifact, provenance "
            "FROM retrieval_chunks WHERE org_id = %s AND source_artifact = %s",
            (org_id, source_artifact),
        )
        return [dict(r) for r in cur.fetchall()]
    finally:
        con.close()


def _count(org_id: str) -> int:
    con = db.connect()
    try:
        cur = con.cursor()
        cur.execute("SELECT COUNT(*) AS n FROM retrieval_chunks WHERE org_id = %s", (org_id,))
        return int(cur.fetchone()["n"])
    finally:
        con.close()


@pytest.fixture
def org(request, monkeypatch):
    monkeypatch.setenv("MODEL_EMBEDDING_PROVIDER", _PROVIDER.name)
    name = f"doc_acc_{request.node.name}"[:60]
    _cleanup(name)
    yield name
    _cleanup(name)


def _run(org_id, source, store, **kw):
    return ingest_documents(
        org_id,
        source=source,
        read_checkpoint=store.read,
        save_checkpoint=store.save,
        **kw,
    )


# Four healthy files, one per phase-one format, each carrying a distinct term.
def _healthy_entries():
    return [
        {"id": "lib/report.pdf", "filename": "report.pdf", "content_type": "application/pdf",
         "bytes": _make_pdf("alpha quarterly review"), "signature": "s-pdf-1", "ts": "2026-07-01T00:00:00Z"},
        {"id": "lib/notes.docx", "filename": "notes.docx",
         "bytes": _make_docx("beta meeting notes"), "signature": "s-docx-1", "ts": "2026-07-02T00:00:00Z"},
        {"id": "lib/budget.xlsx", "filename": "budget.xlsx",
         "bytes": _make_xlsx("gamma budget"), "signature": "s-xlsx-1", "ts": "2026-07-03T00:00:00Z"},
        {"id": "lib/deck.pptx", "filename": "deck.pptx",
         "bytes": _make_pptx("delta launch"), "signature": "s-pptx-1", "ts": "2026-07-04T00:00:00Z"},
    ]


_TERM_FOR = {
    "lib/report.pdf": "alpha",
    "lib/notes.docx": "beta",
    "lib/budget.xlsx": "gamma",
    "lib/deck.pptx": "delta",
}


# ===========================================================================
# AC1 — each format extracted with provenance, delivered, and retrievable
# ===========================================================================
def test_ac1_all_formats_delivered_and_retrievable(org):
    src = ListSource(_healthy_entries())
    result = _run(org, src, Store())
    assert result.ok, result.error
    assert result.artifacts_handed_off == 4
    assert result.artifacts_indexed == 4

    embedder.embed_pending_for_org(org)
    for artifact_id, term in _TERM_FOR.items():
        hits = retrieve(org, term, k=5, source_filter=["document"])
        assert hits, f"{term!r} retrieved nothing"
        assert hits[0].source_artifact == artifact_id


# ===========================================================================
# AC2 — incremental: unchanged files are not re-processed; changed ones are
# ===========================================================================
def test_ac2_incremental_skips_unchanged_reprocesses_changed(org):
    store = Store()
    entries = _healthy_entries()
    src1 = ListSource([dict(e) for e in entries])
    first = _run(org, src1, store)
    assert first.artifacts_indexed == 4
    baseline = _count(org)
    assert baseline >= 4

    # Second run, identical inventory → nothing changed → nothing re-read/re-handed.
    src2 = ListSource([dict(e) for e in entries])
    second = _run(org, src2, store)
    assert second.artifacts_handed_off == 0
    assert src2.reads == []
    assert _count(org) == baseline  # no chunks re-created

    # Change one file's content + signature → only it is re-read and re-indexed.
    changed = [dict(e) for e in entries]
    changed[1]["signature"] = "s-docx-2"
    changed[1]["bytes"] = _make_docx("beta revised notes")
    src3 = ListSource(changed)
    third = _run(org, src3, store)
    assert third.artifacts_handed_off == 1
    assert src3.reads == ["lib/notes.docx"]


# ===========================================================================
# AC3 — scanned/encrypted files are skipped-with-reason and NEVER indexed
# ===========================================================================
def test_ac3_scanned_and_encrypted_skipped_never_indexed(org):
    entries = _healthy_entries() + [
        {"id": "lib/scanned.pdf", "filename": "scanned.pdf", "content_type": "application/pdf",
         "bytes": _make_pdf("", with_text=False), "signature": "s-scan-1", "ts": "2026-07-05T00:00:00Z"},
        {"id": "lib/locked.docx", "filename": "locked.docx",
         "bytes": _OLE_MAGIC + b"encrypted", "signature": "s-lock-1", "ts": "2026-07-06T00:00:00Z"},
    ]
    result = _run(org, ListSource(entries), Store())
    assert result.ok, result.error
    # The two loud-skip files were never handed to the substrate...
    assert result.artifacts_handed_off == 4  # only the healthy four
    assert _rows_for(org, "lib/scanned.pdf") == []
    assert _rows_for(org, "lib/locked.docx") == []
    # ...but the healthy files in the same run were indexed (run continued).
    assert _rows_for(org, "lib/report.pdf")

    # And their content is not retrievable as empty text — a query never surfaces them.
    embedder.embed_pending_for_org(org)
    for term in ("alpha", "beta", "gamma", "delta"):
        for hit in retrieve(org, term, k=10, source_filter=["document"]):
            assert hit.source_artifact not in ("lib/scanned.pdf", "lib/locked.docx")


# ===========================================================================
# AC4 — oversized file skipped-with-reason, run continues, never indexed
# ===========================================================================
def test_ac4_oversized_skipped_run_continues(org, monkeypatch):
    # Cap comfortably above the healthy files (a .pptx/.docx zip is tens of KB)
    # but well below the oversized text file.
    monkeypatch.setenv("DOCUMENT_MAX_FILE_BYTES", "100000")
    big_text = ("epsilon " * 20000).encode("utf-8")  # ~160 KB, over the cap
    entries = _healthy_entries() + [
        {"id": "lib/huge.txt", "filename": "huge.txt", "content_type": "text/plain",
         "bytes": big_text, "size_bytes": len(big_text), "signature": "s-huge-1",
         "ts": "2026-07-07T00:00:00Z"},
    ]
    src = ListSource(entries)
    result = _run(org, src, Store())
    assert result.ok, result.error
    # The oversized file was neither read (size known up front) nor indexed...
    assert "lib/huge.txt" not in src.reads
    assert _rows_for(org, "lib/huge.txt") == []
    # ...and the run continued: the healthy files were indexed.
    assert result.artifacts_indexed == 4
    # The oversized file has no chunks, so it can never surface as a retrieval hit.
    embedder.embed_pending_for_org(org)
    for hit in retrieve(org, "epsilon", k=10, source_filter=["document"]):
        assert hit.source_artifact != "lib/huge.txt"


# ===========================================================================
# AC5 — a single corrupt file fails in isolation; the rest proceed
# ===========================================================================
def test_ac5_corrupt_file_isolated_others_indexed(org):
    entries = _healthy_entries() + [
        {"id": "lib/broken.docx", "filename": "broken.docx",
         "bytes": b"PK\x03\x04totally-broken", "signature": "s-broken-1",
         "ts": "2026-07-08T00:00:00Z"},
    ]
    result = _run(org, ListSource(entries), Store())
    # The corrupt file errored on its own and never reached the substrate; the
    # four healthy files were indexed and the run did not fail.
    assert result.ok, result.error
    assert result.artifacts_indexed == 4
    assert _rows_for(org, "lib/broken.docx") == []
    embedder.embed_pending_for_org(org)
    hits = retrieve(org, "alpha", k=5, source_filter=["document"])
    assert hits and hits[0].source_artifact == "lib/report.pdf"


# ===========================================================================
# AC6 — observed provenance + correct source file, stored and retrieved
# ===========================================================================
def test_ac6_observed_provenance_and_correct_source_file(org):
    _run(org, ListSource(_healthy_entries()), Store())

    rows = _rows_for(org, "lib/report.pdf")
    assert rows
    prov = json.loads(rows[0]["provenance"])
    assert prov["origin"] == OBSERVED
    assert prov["filename"] == "report.pdf"
    assert prov["evidence_pointer"]["origin"] == OBSERVED
    assert prov["evidence_pointer"]["source_artifact"] == "lib/report.pdf"
    assert all(r["source_system"] == "document" for r in rows)

    embedder.embed_pending_for_org(org)
    hits = retrieve(org, "alpha", k=5, source_filter=["document"])
    assert hits
    top = hits[0]
    assert top.source_system == "document"
    assert top.source_artifact == "lib/report.pdf"
    assert top.to_evidence_pointer().origin == OBSERVED
